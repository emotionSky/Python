from pptx import Presentation  # pip install python-pptx --user
from docx import Document  # pip install python-docx --user
from pathlib import Path

dst_folder = Path('./输出文件')
if not dst_folder.exists():
    dst_folder.mkdir(parents=True)

word_file = Document()
file_path = './素材/员工管理制度.pptx'
ppt = Presentation(file_path)
for i in ppt.slides:  # 遍历演示文稿中的所有幻灯片
    for j in i.shapes:  # 遍历幻灯片中的所有形状
        if j.has_text_frame:  # 判断是否含有文本框
            text_frame = j.text_frame
            for paragraph in text_frame.paragraphs:
                word_file.add_paragraph(paragraph.text)

word_file.save(dst_folder / '员工管理制度.docx')
