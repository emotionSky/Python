from pathlib import Path
from openpyxl import load_workbook  # pip install openpyxl --user
from pptx import Presentation  # pip install python-pptx --user
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

src_folder = Path('./素材/公司产品介绍')
dst_folder = Path('./输出文件')
if not dst_folder.exists():
    dst_folder.mkdir(parents=True)

ppt_file = Presentation()
# Cm() 函数的数值单位就是 cm
sl_w = Cm(40)                 # 设置幻灯片宽度
sl_h = Cm(22.5)               # 设置幻灯片高度
ppt_file.slide_width = sl_w   # 定义幻灯片宽度
ppt_file.slide_height = sl_h  # 定义幻灯片高度
wb = load_workbook(src_folder / '产品信息表.xlsx')
ws = wb.active  # 激活工作表

for row in range(2, ws.max_row + 1):
    # 插入一张不带任何占位符的空白幻灯片，这里的 slide_layouts 就是PPT新建一张幻灯片的时候可供选择的样式
    slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[6])
    # 设置幻灯片的背景为纯色填充
    slide.background.fill.solid()
    # 设置幻灯片的背景填充颜色
    slide.background.fill.fore_color.rgb = RGBColor(230, 230, 230)

    # 插入形状
    # 在幻灯片添加形状，后插入的内容会覆盖先插入的内容
    rec = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,       # 矩形
        left=0,                    # 形状的左上角到幻灯片的左侧距离
        top=(sl_h - Cm(10)) // 2,  # 形状的左上角到幻灯片顶部的距离
        width=sl_w,                # 形状的宽度
        height=Cm(10))             # 形状的高度
    
    rec.fill.solid()                                 # 形状的背景为纯色填充
    rec.fill.fore_color.rgb = RGBColor(197, 90, 17)  # 形状的填充颜色为深橙色
    rec.line.fill.background()                       # 形状的轮廓为无

    # 插入图片
    # 这里需要将产品和对应的图片联系起来，不能插错了
    pic_path = src_folder / (str(ws.cell(row=row, column=1).value) + '.png')
    # 这里省略了设置图片高度的参数，表示按照原始宽高比例自动计算
    pic = slide.shapes.add_picture(
        str(pic_path),  # 图片路径
        left=Cm(2),     # 图片左上角到幻灯片左侧的距离
        top=0,          # 图片左上角到幻灯片顶部的距离
        width=Cm(15))   # 图片的宽度
    # 重新设置图片和幻灯片顶部的距离
    pic.top = (sl_h - pic.height) // 2
    # 为图片设置跳转的网页
    pic.click_action.hyperlink.address = str(ws.cell(row=row, column=5).value)

    # 插入文本
    txt = slide.shapes.add_textbox(
        left=sl_w / 2,   # 文本框左上角到幻灯片左侧的距离
        top=0,           # 文本框左上角到幻灯片顶部的距离
        width=sl_w / 2,  # 文本框的宽度
        height=sl_h)     # 文本框的高度
    # 文本框的垂直对齐方式： 居中
    txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # 文本框中的文本超过文本宽度时自动换行
    txt.text_frame.word_wrap = True
    # 依次读取工作表中的 2-4 列
    for col in range(2, 5):
        p = txt.text_frame.add_paragraph()  # 添加段落
        # 段落内容
        p.text = f'【{ws.cell(row=1, column=col).value}】{ws.cell(row=row, column=col).value}\n'
        # 文字的对齐方式
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        p.font.name = '华文中宋'                        # 字体
        p.font.size = Pt(28)                        # 字号
        p.font.color.rgb = RGBColor(255, 255, 255)  # 文字颜色

ppt_file.save(dst_folder / '公司产品介绍.pptx')
